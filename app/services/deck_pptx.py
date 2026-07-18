"""Business analysis deck (PPTX) with real matplotlib charts."""

from __future__ import annotations

import logging
from datetime import datetime
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from app.config import get_settings

logger = logging.getLogger(__name__)

ACCENT = RGBColor(0x1F, 0x4E, 0x79)


def generate_analysis_deck(
    data: dict[str, Any],
    prefs: dict[str, str],
    low_stock: list[dict[str, Any]] | None = None,
) -> str:
    """data = analytics_service.sales_for_range(...) dict. Returns pptx path."""
    settings = get_settings()
    out_dir = Path(settings.generated_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    charts_dir = out_dir / f"charts_{stamp}"
    charts_dir.mkdir(parents=True, exist_ok=True)

    shop_name = prefs.get("shop_name", "Kirana Store")
    low_stock = low_stock or []

    trend_png = _chart_sales_trend(data.get("daily", []), charts_dir)
    top_png = _chart_top_items(data.get("top_items", []), charts_dir)
    pay_png = _chart_payment_split(data, charts_dir)
    stock_png = _chart_stock_health(low_stock, charts_dir)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    _title_slide(
        prs,
        f"{shop_name} — Sales Analysis",
        f"{data.get('start_date')} to {data.get('end_date')}",
    )

    # KPI slide
    _kpi_slide(prs, data)

    # Chart slides
    if trend_png:
        _image_slide(prs, "Daily Sales Trend", trend_png)
    if top_png:
        _image_slide(prs, "Top Selling Items", top_png)
    if pay_png:
        _image_slide(prs, "Payment Mode Split", pay_png)
    if stock_png:
        _image_slide(prs, "Stock Health (Low / Reorder)", stock_png)

    _takeaways_slide(prs, data, low_stock)

    file_path = str(out_dir / f"analysis_{data.get('start_date')}_{data.get('end_date')}_{stamp}.pptx")
    prs.save(file_path)
    logger.info("analysis_deck_generated path=%s", file_path)
    return file_path


# --------------------------- charts ---------------------------------------

def _chart_sales_trend(daily: list[dict[str, Any]], out: Path) -> str | None:
    if not daily:
        return None
    dates = [d["date"][5:] for d in daily]
    sales = [d["total_sales"] for d in daily]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.plot(dates, sales, marker="o", color="#1f4e79", linewidth=2)
    ax.fill_between(range(len(sales)), sales, alpha=0.1, color="#1f4e79")
    ax.set_ylabel("Sales (₹)")
    ax.set_title("Sales by Day")
    ax.grid(True, alpha=0.3)
    fig.tight_layout()
    p = str(out / "trend.png")
    fig.savefig(p, dpi=130)
    plt.close(fig)
    return p


def _chart_top_items(top: list[dict[str, Any]], out: Path) -> str | None:
    if not top:
        return None
    names = [t["sku_name"] for t in top][::-1]
    amounts = [t["amount"] for t in top][::-1]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.barh(names, amounts, color="#2e75b6")
    ax.set_xlabel("Revenue (₹)")
    ax.set_title("Top Items by Revenue")
    for i, v in enumerate(amounts):
        ax.text(v, i, f" ₹{v:,.0f}", va="center", fontsize=8)
    fig.tight_layout()
    p = str(out / "top.png")
    fig.savefig(p, dpi=130)
    plt.close(fig)
    return p


def _chart_payment_split(data: dict[str, Any], out: Path) -> str | None:
    modes = {
        "Cash": data.get("cash_total", 0),
        "UPI": data.get("upi_total", 0),
        "Card": data.get("card_total", 0),
        "Khata": data.get("khata_total", 0),
    }
    modes = {k: v for k, v in modes.items() if v > 0}
    if not modes:
        return None
    fig, ax = plt.subplots(figsize=(7, 4.5))
    ax.pie(
        list(modes.values()),
        labels=list(modes.keys()),
        autopct="%1.1f%%",
        colors=["#1f4e79", "#2e75b6", "#9dc3e6", "#c55a11"],
        startangle=90,
    )
    ax.set_title("Payment Mode Split")
    fig.tight_layout()
    p = str(out / "pay.png")
    fig.savefig(p, dpi=130)
    plt.close(fig)
    return p


def _chart_stock_health(low_stock: list[dict[str, Any]], out: Path) -> str | None:
    if not low_stock:
        return None
    items = low_stock[:10]
    names = [i["sku_name"] for i in items][::-1]
    qty = [i["quantity_on_hand"] for i in items][::-1]
    reorder = [i["reorder_level"] for i in items][::-1]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    y = range(len(names))
    ax.barh(y, qty, color="#c00000", label="On hand")
    ax.plot(reorder, y, "D", color="#000000", label="Reorder level")
    ax.set_yticks(list(y))
    ax.set_yticklabels(names)
    ax.set_xlabel("Quantity")
    ax.set_title("Low Stock Items")
    ax.legend()
    fig.tight_layout()
    p = str(out / "stock.png")
    fig.savefig(p, dpi=130)
    plt.close(fig)
    return p


# --------------------------- slides ----------------------------------------

def _title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def _kpi_slide(prs: Presentation, data: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, "Key Numbers")
    kpis = [
        ("Total Sales", f"₹{data.get('total_sales', 0):,.2f}"),
        ("GST Collected", f"₹{data.get('tax_collected', 0):,.2f}"),
        ("Bills", str(data.get("bills_count", 0))),
        ("UPI", f"₹{data.get('upi_total', 0):,.2f}"),
        ("Cash", f"₹{data.get('cash_total', 0):,.2f}"),
        ("Khata", f"₹{data.get('khata_total', 0):,.2f}"),
    ]
    left = 0.8
    top = 1.8
    for i, (label, value) in enumerate(kpis):
        col = i % 3
        rowi = i // 3
        box = slide.shapes.add_textbox(
            Inches(left + col * 4.1), Inches(top + rowi * 2.1), Inches(3.8), Inches(1.8)
        )
        tf = box.text_frame
        tf.text = value
        tf.paragraphs[0].font.size = Pt(30)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = ACCENT
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def _image_slide(prs: Presentation, title: str, img_path: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, title)
    slide.shapes.add_picture(img_path, Inches(1.4), Inches(1.4), height=Inches(5.4))


def _takeaways_slide(
    prs: Presentation, data: dict[str, Any], low_stock: list[dict[str, Any]]
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, "Takeaways")
    lines = []
    top = data.get("top_items", [])
    if top:
        lines.append(f"Best seller: {top[0]['sku_name']} (₹{top[0]['amount']:,.0f})")
    lines.append(f"GST to remit: ₹{data.get('tax_collected', 0):,.2f}")
    if data.get("upi_total", 0) > data.get("cash_total", 0):
        lines.append("UPI is the dominant payment mode.")
    if low_stock:
        names = ", ".join(i["sku_name"] for i in low_stock[:5])
        lines.append(f"Reorder soon: {names}")
    else:
        lines.append("Stock levels healthy — nothing below reorder.")

    box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + line
        p.font.size = Pt(20)
        p.space_after = Pt(10)


def _heading(slide: Any, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT
